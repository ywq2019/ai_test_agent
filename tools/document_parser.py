"""
文档解析工具 - 支持 PDF / DOCX / TXT / MD / CSV / XLSX / HTML / PPTX / JSON
"""
import re
import csv
import json
from io import StringIO
from pathlib import Path
from typing import Dict, List, Any, Optional
from loguru import logger

# ── 可选依赖，缺失时优雅降级 ──────────────────────────────────────────
# pymupdf (fitz) 优先，中文 PDF 兼容性更好；不可用时降级到 PyPDF2
try:
    import fitz as _fitz  # pymupdf
    _HAS_FITZ = True
except ImportError:
    _fitz = None
    _HAS_FITZ = False

try:
    from PyPDF2 import PdfReader
except ImportError:
    PdfReader = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from html.parser import HTMLParser as _HTMLParser
    _HAS_HTML_PARSER = True
except ImportError:
    _HAS_HTML_PARSER = False


# ── 轻量 HTML → text 提取器 ──────────────────────────────────────────
class _TextExtractor(_HTMLParser if _HAS_HTML_PARSER else object):
    def __init__(self):
        if _HAS_HTML_PARSER:
            super().__init__()
        self._parts: List[str] = []
        self._skip = False

    def handle_starttag(self, tag, attrs):
        self._skip = tag in ("script", "style")

    def handle_endtag(self, tag):
        if tag in ("script", "style"):
            self._skip = False

    def handle_data(self, data):
        if not self._skip and data.strip():
            self._parts.append(data.strip())

    def get_text(self) -> str:
        return "\n".join(self._parts)


def _detect_encoding(raw: bytes) -> str:
    """简单编码探测：utf-8 → gbk → latin-1"""
    for enc in ("utf-8", "gbk", "utf-16", "latin-1"):
        try:
            raw.decode(enc)
            return enc
        except (UnicodeDecodeError, LookupError):
            pass
    return "latin-1"


class DocumentParser:
    SUPPORTED = {
        ".pdf", ".docx", ".doc",
        ".txt", ".md",
        ".csv",
        ".xlsx", ".xls",
        ".html", ".htm",
        ".pptx",
        ".json",
    }

    def is_supported(self, file_path: str) -> bool:
        return Path(file_path).suffix.lower() in self.SUPPORTED

    async def parse(self, file_path: str) -> Dict[str, Any]:
        ext = Path(file_path).suffix.lower()
        parsers = {
            ".pdf":  self._parse_pdf,
            ".docx": self._parse_docx,
            ".doc":  self._parse_docx,
            ".txt":  self._parse_txt,
            ".md":   self._parse_md,
            ".csv":  self._parse_csv,
            ".xlsx": self._parse_xlsx,
            ".xls":  self._parse_xlsx,
            ".html": self._parse_html,
            ".htm":  self._parse_html,
            ".pptx": self._parse_pptx,
            ".json": self._parse_json,
        }
        fn = parsers.get(ext)
        if fn is None:
            raise ValueError(f"不支持的文档格式: {ext}，支持格式: {', '.join(sorted(self.SUPPORTED))}")
        return await fn(file_path)

    # ── PDF ──────────────────────────────────────────────────────────
    async def _parse_pdf(self, file_path: str) -> Dict[str, Any]:
        """
        PDF 文本提取：优先使用 pymupdf(fitz)，中文兼容性更好；
        不可用时降级到 PyPDF2。
        """
        pages_text = []

        if _HAS_FITZ:
            # pymupdf — 正确处理中文 PDF（Unicode 直接提取）
            doc = _fitz.open(file_path)
            for i, page in enumerate(doc):
                t = page.get_text() or ""
                pages_text.append({"page": i + 1, "text": t})
            doc.close()
            logger.info(f"PDF 解析完成 (pymupdf): {len(pages_text)} 页")
        elif PdfReader is not None:
            # PyPDF2 降级（部分中文 PDF 可能乱码）
            reader = PdfReader(file_path)
            for i, page in enumerate(reader.pages):
                t = page.extract_text() or ""
                pages_text.append({"page": i + 1, "text": t})
            logger.info(f"PDF 解析完成 (PyPDF2): {len(pages_text)} 页")
        else:
            raise ImportError("请安装 PDF 解析库：pip install pymupdf")

        content = "\n".join(p["text"] for p in pages_text)
        return {
            "content": content,
            "page_count": len(pages_text),
            "paragraph_count": len([l for l in content.splitlines() if l.strip()]),
            "pages": pages_text,
            "structured": self._structure_content(content),
            "metadata": {"format": "pdf", "pages": len(pages_text)},
        }

    # ── DOCX ─────────────────────────────────────────────────────────
    async def _parse_docx(self, file_path: str) -> Dict[str, Any]:
        if DocxDocument is None:
            raise ImportError("请安装 python-docx：pip install python-docx")
        doc = DocxDocument(file_path)
        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        tables = []
        for tbl in doc.tables:
            tables.append([[cell.text for cell in row.cells] for row in tbl.rows])
        content = "\n".join(paras)
        logger.info(f"DOCX 解析完成: {len(paras)} 段落, {len(tables)} 表格")
        return {
            "content": content,
            "page_count": 0,
            "paragraph_count": len(paras),
            "table_count": len(tables),
            "tables": tables,
            "structured": self._structure_content(content),
            "metadata": {"format": "docx", "paragraphs": len(paras), "tables": len(tables)},
        }

    # ── TXT ──────────────────────────────────────────────────────────
    async def _parse_txt(self, file_path: str) -> Dict[str, Any]:
        raw = Path(file_path).read_bytes()
        enc = _detect_encoding(raw)
        content = raw.decode(enc, errors="replace")
        lines = [l for l in content.splitlines() if l.strip()]
        logger.info(f"TXT 解析完成: {len(lines)} 行, 编码 {enc}")
        return {
            "content": content,
            "page_count": 0,
            "paragraph_count": len(lines),
            "structured": self._structure_content(content),
            "metadata": {"format": "txt", "encoding": enc, "lines": len(lines)},
        }

    # ── Markdown ─────────────────────────────────────────────────────
    async def _parse_md(self, file_path: str) -> Dict[str, Any]:
        raw = Path(file_path).read_bytes()
        enc = _detect_encoding(raw)
        content = raw.decode(enc, errors="replace")
        # 去掉 markdown 语法符号，方便后续提取
        plain = re.sub(r"#{1,6}\s", "", content)
        plain = re.sub(r"\*{1,2}(.+?)\*{1,2}", r"\1", plain)
        plain = re.sub(r"`{1,3}[^`]*`{1,3}", "", plain)
        plain = re.sub(r"\[(.+?)\]\(.+?\)", r"\1", plain)
        lines = [l for l in plain.splitlines() if l.strip()]
        headings = [l.strip() for l in content.splitlines() if re.match(r"^#{1,6}\s", l)]
        logger.info(f"Markdown 解析完成: {len(headings)} 个标题")
        return {
            "content": content,
            "page_count": 0,
            "paragraph_count": len(lines),
            "headings": headings,
            "structured": self._structure_content(plain),
            "metadata": {"format": "md", "encoding": enc, "headings": len(headings)},
        }

    # ── CSV ──────────────────────────────────────────────────────────
    async def _parse_csv(self, file_path: str) -> Dict[str, Any]:
        raw = Path(file_path).read_bytes()
        enc = _detect_encoding(raw)
        text = raw.decode(enc, errors="replace")
        reader = csv.reader(StringIO(text))
        rows = list(reader)
        headers = rows[0] if rows else []
        # 将表格转换为易读文本
        lines = [", ".join(r) for r in rows]
        content = "\n".join(lines)
        logger.info(f"CSV 解析完成: {len(rows)} 行, {len(headers)} 列")
        return {
            "content": content,
            "page_count": 0,
            "paragraph_count": len(rows),
            "headers": headers,
            "row_count": len(rows),
            "col_count": len(headers),
            "structured": self._structure_content(content),
            "metadata": {"format": "csv", "rows": len(rows), "columns": len(headers)},
        }

    # ── Excel (xlsx/xls) ─────────────────────────────────────────────
    async def _parse_xlsx(self, file_path: str) -> Dict[str, Any]:
        if openpyxl is None:
            raise ImportError("请安装 openpyxl：pip install openpyxl")
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        all_lines = []
        sheet_info = []
        for sheet in wb.worksheets:
            rows_text = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    rows_text.append(", ".join(cells))
            all_lines.append(f"[Sheet: {sheet.title}]")
            all_lines.extend(rows_text)
            sheet_info.append({"name": sheet.title, "rows": len(rows_text)})
        wb.close()
        content = "\n".join(all_lines)
        total_rows = sum(s["rows"] for s in sheet_info)
        logger.info(f"Excel 解析完成: {len(sheet_info)} 个 Sheet, {total_rows} 行")
        return {
            "content": content,
            "page_count": len(sheet_info),
            "paragraph_count": total_rows,
            "sheets": sheet_info,
            "structured": self._structure_content(content),
            "metadata": {"format": "xlsx", "sheets": len(sheet_info), "total_rows": total_rows},
        }

    # ── HTML ─────────────────────────────────────────────────────────
    async def _parse_html(self, file_path: str) -> Dict[str, Any]:
        raw = Path(file_path).read_bytes()
        enc = _detect_encoding(raw)
        html_text = raw.decode(enc, errors="replace")

        # 优先用 BeautifulSoup 精准提取主内容区域
        content = self._extract_html_main_content(html_text)

        # 如果主内容提取结果太少（SPA/数据导出类 HTML），尝试从 script 标签里提取 JSON 数据
        if len(content) < 200:
            content = self._extract_script_data(html_text) or content

        # BeautifulSoup 不可用时，退回轻量提取器
        if not content:
            if _HAS_HTML_PARSER:
                extractor = _TextExtractor()
                extractor.feed(html_text)
                content = extractor.get_text()
            else:
                content = re.sub(r"<[^>]+>", " ", html_text)
            content = re.sub(r"\s{3,}", "\n\n", content).strip()

        lines = [l for l in content.splitlines() if l.strip()]
        logger.info(f"HTML 解析完成: {len(lines)} 行")
        return {
            "content": content,
            "page_count": 0,
            "paragraph_count": len(lines),
            "structured": self._structure_content(content),
            "metadata": {"format": "html", "encoding": enc},
        }

    def _extract_html_main_content(self, html_text: str) -> str:
        """用 BeautifulSoup 清洗 HTML，提取主需求内容区域。
        策略（保守优先，避免误删正文）：
        1. 只删语义明确的噪音标签（script/style/nav/footer 等 HTML5 结构标签）
        2. 优先从语义主内容容器提取，找到就只取那一块
        3. 找不到则全文 get_text，再做行级过滤去掉明显 UI 碎片
        """
        try:
            from bs4 import BeautifulSoup, Comment

            soup = BeautifulSoup(html_text, "html.parser")

            # Step1: 只删语义明确的噪音标签（不做 class 关键词匹配，避免误伤正文）
            _NOISE_TAGS = [
                "script", "style", "head", "meta", "link", "noscript",
                "nav", "footer", "aside", "iframe", "svg", "canvas",
            ]
            for tag in soup(_NOISE_TAGS):
                tag.decompose()
            # 删 HTML 注释
            for comment in soup.find_all(string=lambda s: isinstance(s, Comment)):
                comment.extract()

            # Step2: 优先从语义主内容区提取（找到就只取那一块，大幅减少噪音）
            _MAIN_SELECTORS = [
                "main", "article",
                "[role='main']",
                "#content", "#main", "#main-content", "#page-content",
                "#app-content", "#editor-content",
                ".main-content", ".page-content", ".article-content",
                ".wiki-content", ".doc-content", ".markdown-body",
                ".ql-editor", ".ProseMirror",           # 富文本编辑器
                ".requirement-content", ".spec-content",
            ]
            raw_text = ""
            for sel in _MAIN_SELECTORS:
                try:
                    node = soup.select_one(sel)
                    if node:
                        candidate = node.get_text(separator="\n", strip=True)
                        if len(candidate) >= 200:
                            raw_text = candidate
                            logger.info(f"HTML 主内容区命中: {sel}，{len(raw_text)} 字")
                            break
                except Exception:
                    continue

            # 未命中主内容区，退回全文 get_text
            if not raw_text:
                raw_text = soup.get_text(separator="\n", strip=True)
                logger.info(f"HTML 未命中主内容区，使用全文 get_text，{len(raw_text)} 字")

            # Step3: 行级过滤——只去掉极短碎片和纯符号行，保留所有正文
            lines_kept = []
            for line in raw_text.splitlines():
                line = line.strip()
                if not line:
                    continue
                # 过滤：纯空白 / 纯数字符号行（如页码、分隔线）
                if re.fullmatch(r'[\d\s\-–—|/\\·•○●□■◆▷▶→←★☆※…。，,\.]+', line):
                    continue
                # 过滤：极短行（3字以内），通常是图标/数字/单字按钮
                if len(line) <= 3:
                    continue
                lines_kept.append(line)

            result = re.sub(r'\n{3,}', '\n\n', "\n".join(lines_kept)).strip()
            logger.info(f"HTML 主内容提取: {len(html_text)} → {len(result)} 字")
            return result

        except ImportError:
            logger.debug("BeautifulSoup 不可用，跳过主内容提取")
            return ""
        except Exception as e:
            logger.warning(f"HTML 主内容提取失败: {e}")
            return ""

    def _extract_script_data(self, html_text: str) -> str:
        """从 script 标签中提取 JS 变量赋值的 JSON 数据，转为可读文本。"""
        # 用字符串搜索代替正则，避免大文件回溯爆炸
        pos = 0
        while True:
            start = html_text.find("<script", pos)
            if start == -1:
                break
            tag_end = html_text.find(">", start)
            if tag_end == -1:
                break
            content_start = tag_end + 1
            end = html_text.find("</script>", content_start)
            if end == -1:
                end = len(html_text)
            script = html_text[content_start:end].strip()
            pos = end + 9

            if len(script) < 100:
                continue
            # 匹配 const/let/var xxx = [...] 或 = {...}
            m = re.match(r"(?:const|let|var)\s+\w+\s*=\s*(\[|{)", script)
            if not m:
                continue
            json_start_idx = m.start(1)
            json_str = script[json_start_idx:]
            try:
                data, _ = json.JSONDecoder().raw_decode(json_str)
                result = self._flatten_json_to_text(data)
                if len(result) > 100:
                    logger.info(f"从 script 标签提取文本: {len(result)} 字符")
                    return result
            except (json.JSONDecodeError, ValueError):
                pass
        return ""

    def _flatten_json_to_text(self, data, depth: int = 0, _acc: list = None) -> str:
        """递归把 JSON 数据转成纯文本，方便 AI 理解。超过 30000 字符时提前停止。"""
        if _acc is None:
            _acc = []
        indent = "  " * depth
        if isinstance(data, list):
            for item in data:
                if sum(len(x) for x in _acc) >= 30000:
                    break
                self._flatten_json_to_text(item, depth, _acc)
        elif isinstance(data, dict):
            for k, v in data.items():
                if sum(len(x) for x in _acc) >= 30000:
                    break
                if isinstance(v, (dict, list)):
                    _acc.append(f"{indent}{k}:\n")
                    self._flatten_json_to_text(v, depth + 1, _acc)
                elif v and str(v).strip():
                    _acc.append(f"{indent}{k}: {v}\n")
        else:
            if data and str(data).strip():
                _acc.append(f"{indent}{data}\n")
        return "".join(_acc)

    # ── PowerPoint ───────────────────────────────────────────────────
    async def _parse_pptx(self, file_path: str) -> Dict[str, Any]:
        if Presentation is None:
            raise ImportError("请安装 python-pptx：pip install python-pptx")
        prs = Presentation(file_path)
        slides_text = []
        for i, slide in enumerate(prs.slides):
            parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
            slides_text.append({"slide": i + 1, "text": "\n".join(parts)})
        content = "\n\n".join(f"[幻灯片 {s['slide']}]\n{s['text']}" for s in slides_text)
        logger.info(f"PPTX 解析完成: {len(slides_text)} 页")
        return {
            "content": content,
            "page_count": len(slides_text),
            "paragraph_count": len([l for l in content.splitlines() if l.strip()]),
            "slides": slides_text,
            "structured": self._structure_content(content),
            "metadata": {"format": "pptx", "slides": len(slides_text)},
        }

    # ── JSON ─────────────────────────────────────────────────────────
    async def _parse_json(self, file_path: str) -> Dict[str, Any]:
        raw = Path(file_path).read_bytes()
        enc = _detect_encoding(raw)
        text = raw.decode(enc, errors="replace")
        try:
            data = json.loads(text)
            content = json.dumps(data, ensure_ascii=False, indent=2)
        except json.JSONDecodeError:
            content = text
            data = {}
        lines = [l for l in content.splitlines() if l.strip()]
        logger.info(f"JSON 解析完成: {len(lines)} 行")
        return {
            "content": content,
            "page_count": 0,
            "paragraph_count": len(lines),
            "structured": self._structure_content(content),
            "metadata": {"format": "json"},
        }

    # ── 通用结构化提取 ────────────────────────────────────────────────
    def _structure_content(self, text: str) -> Dict[str, Any]:
        lines = [l.strip() for l in text.splitlines() if l.strip()]
        sections, current_title, items = [], None, []
        for line in lines:
            if self._is_heading(line):
                if current_title:
                    sections.append({"title": current_title, "items": items})
                current_title, items = line, []
            else:
                items.append(line)
        if current_title:
            sections.append({"title": current_title, "items": items})
        return {
            "sections": sections,
            "functional_points": self._extract_functional_points(text),
            "validation_rules": self._extract_validation_rules(text),
            "business_flows": self._extract_business_flows(text),
        }

    def _is_heading(self, line: str) -> bool:
        patterns = [
            r"^\d+[\.、]", r"^[一二三四五六七八九十]+[\.、]",
            r"^[A-Z][\.\)]", r"^\[.+\]$", r"^#{1,6}\s",
        ]
        return any(re.match(p, line) for p in patterns) or (len(line) < 50 and line.isupper())

    def _extract_functional_points(self, text: str) -> List[str]:
        pts = []
        for p in [r"功能[：:]\s*(.+?)(?=\n|$)", r"需求[：:]\s*(.+?)(?=\n|$)"]:
            pts.extend(re.findall(p, text))
        return list(set(pts))

    def _extract_validation_rules(self, text: str) -> List[str]:
        rules = []
        for p in [r"校验[规则]?[：:]\s*(.+?)(?=\n|$)", r"规则[：:]\s*(.+?)(?=\n|$)"]:
            rules.extend(re.findall(p, text))
        return list(set(rules))

    def _extract_business_flows(self, text: str) -> List[str]:
        flows = []
        for p in [r"流程[：:]\s*(.+?)(?=\n|$)", r"步骤[：:]\s*(.+?)(?=\n|$)"]:
            flows.extend(re.findall(p, text))
        return list(set(flows))


document_parser = DocumentParser()
